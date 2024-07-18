import json
import sys

from pathlib import Path
from typing import List

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as ppt
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table, _Row, _Cell
from pptx.slide import Slide
from pptx.shapes.group import GroupShape
from werkzeug.datastructures import FileStorage

from magic_doc.contrib.office import OfficeExtractor
from magic_doc.contrib.model import ExtractResponse, Page, Content


class PptxExtractor(OfficeExtractor):
    def __init__(self) -> None:
        super().__init__()

    def setup(self):
        pass

    def handle_shape(
        self,
        shape: Shape,
        content_list: List[Content],
    ):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                content_list.append(
                    Content(
                        type="text",
                        data=paragraph.text + "\n",
                    )
                )
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pass
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shape: GraphicFrame
            table: Table = shape.table
            md = "\n"
            for row_no, row in enumerate(table.rows):
                row: _Row
                md += "|"
                if row_no == 1:
                    for col in row.cells:
                        md += "---|"
                    md += "\n|"
                for col in row.cells:
                    cell: _Cell = col
                    md += " " + cell.text.replace("\r", " ").replace("\n", " ") + " |"
                md += "\n"
            md += "\n"
            content_list.append(Content(type="md", data=md))
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape: GroupShape
            for sub_shape in shape.shapes:
                self.handle_shape(sub_shape, content_list)
        else:
            # print(shape.shape_type, type(shape), file=sys.stderr)
            pass

    def extract(
        self,
        r: FileStorage | Path,
    ) -> ExtractResponse:
        pages = []
   
        presentation: ppt = Presentation(r)
        for page_no, slide in enumerate(presentation.slides):
            slide: Slide
            page = Page(page_no=page_no, content_list=[])
            for shape in slide.shapes:
                self.handle_shape(
                    shape,
                    page["content_list"],
                )

            pages.append(page)
        return pages

